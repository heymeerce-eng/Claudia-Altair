import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ALTAIR brand palette
CREAM      = RGBColor(0xF9, 0xF6, 0xF1)
DARK_BROWN = RGBColor(0x48, 0x38, 0x38)
DARK_RED   = RGBColor(0x47, 0x05, 0x06)
ROSE       = RGBColor(0xA5, 0x8C, 0x8A)
BEIGE      = RGBColor(0xBA, 0xA6, 0x9F)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

DOCS_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "documentos")


def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, text, left, top, width, height,
             size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Arial"
    return tb


def _line(slide, left, top, width, color: RGBColor):
    rect = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        left, top, width, Pt(1),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    return rect


def _cover(prs, title, subtitle, date):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CREAM)
    _line(slide, Inches(1), Inches(0.8), Inches(11.33), DARK_RED)
    _textbox(slide, title,
             Inches(1), Inches(2.2), Inches(11.33), Inches(2.8),
             40, DARK_BROWN, align=PP_ALIGN.CENTER)
    _line(slide, Inches(4.5), Inches(5.0), Inches(4.33), ROSE)
    _textbox(slide, subtitle.upper(),
             Inches(1), Inches(5.3), Inches(11.33), Inches(0.5),
             9, ROSE, align=PP_ALIGN.CENTER)
    if date:
        _textbox(slide, date,
                 Inches(1), Inches(6.0), Inches(11.33), Inches(0.4),
                 9, BEIGE, align=PP_ALIGN.CENTER)
    _line(slide, Inches(1), Inches(6.7), Inches(11.33), DARK_RED)


def _section(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, BEIGE)
    _textbox(slide, title,
             Inches(1.5), Inches(2.8), Inches(10.33), Inches(1.8),
             30, DARK_BROWN, align=PP_ALIGN.CENTER)
    _textbox(slide, "ALTAIR",
             Inches(10.8), Inches(6.9), Inches(1.8), Inches(0.4),
             7, DARK_BROWN, align=PP_ALIGN.RIGHT)


def _content(prs, title, body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CREAM)
    _textbox(slide, title,
             Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.9),
             20, DARK_RED)
    _line(slide, Inches(0.8), Inches(1.45), Inches(11.5), ROSE)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = tb.text_frame
    tf.word_wrap = True

    items = body if isinstance(body, list) else [body]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"— {item}" if isinstance(body, list) else item
        run.font.size = Pt(16)
        run.font.color.rgb = DARK_BROWN
        run.font.name = "Arial"

    _textbox(slide, "ALTAIR",
             Inches(11.5), Inches(6.9), Inches(1.2), Inches(0.3),
             7, BEIGE, align=PP_ALIGN.RIGHT)


def _quote(prs, quote, attribution=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CREAM)
    _textbox(slide, "“",
             Inches(1), Inches(0.6), Inches(2), Inches(1.5),
             72, BEIGE)
    _textbox(slide, quote,
             Inches(1.5), Inches(1.8), Inches(10), Inches(3.5),
             22, DARK_BROWN, align=PP_ALIGN.CENTER, italic=True)
    if attribution:
        _line(slide, Inches(5), Inches(5.5), Inches(3.33), ROSE)
        _textbox(slide, f"— {attribution}",
                 Inches(1.5), Inches(5.8), Inches(10), Inches(0.5),
                 11, ROSE, align=PP_ALIGN.CENTER)


def generate_presentation(
    title: str,
    slides: list,
    subtitle: str = "ALTAIR ACADEMIA",
    date: str = "",
) -> str:
    """
    slides: list of dicts
      - type: "section" | "content" | "quote"
      - title: str (section / content)
      - body: str or list[str] (content)
      - quote: str, attribution: str (quote)
    Returns the saved file path.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _cover(prs, title, subtitle, date)

    for s in slides:
        t = s.get("type", "content")
        if t == "section":
            _section(prs, s.get("title", ""))
        elif t == "quote":
            _quote(prs, s.get("quote", s.get("body", "")), s.get("attribution", ""))
        else:
            _content(prs, s.get("title", ""), s.get("body", ""))

    os.makedirs(DOCS_DIR, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in " -_" else "_" for c in title)
    safe = safe[:40].strip().replace(" ", "_").lower()
    ts = datetime.now().strftime("%Y%m%d_%H%M")
    path = os.path.join(DOCS_DIR, f"presentacion_{safe}_{ts}.pptx")
    prs.save(path)
    return path
